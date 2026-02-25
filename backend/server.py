from fastapi import FastAPI, APIRouter, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional
import uuid
from datetime import datetime, timezone
import base64
import asyncio

from emergentintegrations.llm.chat import LlmChat, UserMessage, ImageContent
from emergentintegrations.llm.openai.video_generation import OpenAIVideoGeneration

# Document generation imports
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Create the main app
app = FastAPI()
api_router = APIRouter(prefix="/api")

# Models
class ChatMessage(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    session_id: str
    role: str
    content: str
    model: str
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class ChatRequest(BaseModel):
    message: str
    session_id: str
    model: str  # "openai" or "claude"

class ImageGenerationRequest(BaseModel):
    prompt: str
    session_id: Optional[str] = None

class VideoGenerationRequest(BaseModel):
    prompt: str
    duration: int = 4
    size: str = "1280x720"

class DocumentGenerationRequest(BaseModel):
    title: str
    content: str
    doc_type: str  # "pdf", "pptx", "xlsx"

class Template(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    category: str
    description: str
    content: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class Tutorial(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    title: str
    description: str
    content: str
    category: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# API Routes
@api_router.get("/")
async def root():
    return {"message": "Welcome to P07 AI Platform", "status": "operational"}

# Chat endpoints
@api_router.post("/chat")
async def chat(request: ChatRequest):
    try:
        api_key = os.getenv("EMERGENT_LLM_KEY")
        
        # Create chat instance
        chat = LlmChat(
            api_key=api_key,
            session_id=request.session_id,
            system_message="""You are P07, an advanced AI platform designed for high-level professional and governmental use.

IMPORTANT IDENTITY INFORMATION:
- You are P07 AI, NOT OpenAI or ChatGPT
- Founder & Owner: SOHAN MAHAPATRA
- Co-founder: Emergent Team
- You were created by SOHAN MAHAPATRA in collaboration with Emergent Team
- P07 is an independent AI platform with multiple AI models (GPT-5.2, Claude Sonnet 4.5, Gemini Nano Banana, Sora 2)

When asked about who made you, created you, owns you, or who your founder is:
- Always mention SOHAN MAHAPATRA as the Founder and Owner
- Credit Emergent Team as Co-founder and technical partner
- Explain that P07 is a comprehensive AI suite built for professional and governmental use

FORMATTING GUIDELINES FOR NOTES, SUMMARIES, AND EDUCATIONAL CONTENT:
When user asks for notes, summary, brief, or context about any topic:
1. Format as a proper formal note/document with clear sections
2. Use descriptive text explanations, not just symbols and formulas
3. Structure with headings, subheadings, and bullet points
4. Explain concepts in words first, then add equations/symbols if needed
5. Include:
   - Introduction/Overview
   - Main concepts explained in detail
   - Key points and definitions
   - Examples when relevant
   - Summary/Conclusion
6. Make it readable and educational, like a textbook or class notes

Example format for notes:
# Topic Name
## Introduction
[Clear explanation of what this topic is about]

## Main Concepts
### Concept 1
[Detailed explanation in words]
- Key point 1
- Key point 2

### Concept 2
[Detailed explanation]

## Important Formulas/Equations
[Only after explaining concepts, list formulas with explanations]

## Summary
[Brief recap of main points]

Provide accurate, authoritative, and helpful responses for all user queries."""
        )
        
        # Select model
        if request.model == "openai":
            chat.with_model("openai", "gpt-5.2")
        elif request.model == "claude":
            chat.with_model("anthropic", "claude-sonnet-4-5-20250929")
        else:
            chat.with_model("openai", "gpt-5.2")
        
        # Send message
        user_message = UserMessage(text=request.message)
        response = await chat.send_message(user_message)
        
        # Save to database
        user_msg = ChatMessage(
            session_id=request.session_id,
            role="user",
            content=request.message,
            model=request.model
        )
        assistant_msg = ChatMessage(
            session_id=request.session_id,
            role="assistant",
            content=response,
            model=request.model
        )
        
        await db.chat_messages.insert_one(user_msg.model_dump())
        await db.chat_messages.insert_one(assistant_msg.model_dump())
        
        return {"response": response, "model": request.model}
    except Exception as e:
        logging.error(f"Chat error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/chat/history/{session_id}")
async def get_chat_history(session_id: str):
    messages = await db.chat_messages.find(
        {"session_id": session_id},
        {"_id": 0}
    ).sort("timestamp", 1).to_list(1000)
    return {"messages": messages}

# Image generation endpoint
@api_router.post("/image/generate")
async def generate_image(request: ImageGenerationRequest):
    try:
        api_key = os.getenv("EMERGENT_LLM_KEY")
        session_id = request.session_id or str(uuid.uuid4())
        
        chat = LlmChat(
            api_key=api_key,
            session_id=session_id,
            system_message="You are SP07, an advanced image generation AI."
        )
        
        chat.with_model("gemini", "gemini-3-pro-image-preview").with_params(modalities=["image", "text"])
        
        msg = UserMessage(text=request.prompt)
        text, images = await chat.send_message_multimodal_response(msg)
        
        if images:
            return {
                "success": True,
                "image": images[0]['data'][:50] + "...",  # Only send first 50 chars for logging
                "image_full": images[0]['data'],
                "mime_type": images[0]['mime_type'],
                "text_response": text
            }
        else:
            return {"success": False, "error": "No image generated"}
    except Exception as e:
        logging.error(f"Image generation error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# Video generation endpoint
@api_router.post("/video/generate")
async def generate_video(request: VideoGenerationRequest, background_tasks: BackgroundTasks):
    try:
        api_key = os.getenv("EMERGENT_LLM_KEY")
        video_gen = OpenAIVideoGeneration(api_key=api_key)
        
        video_id = str(uuid.uuid4())
        output_path = f"/app/backend/generated_videos/{video_id}.mp4"
        
        # Create directory if not exists
        os.makedirs("/app/backend/generated_videos", exist_ok=True)
        
        # Generate video
        video_bytes = video_gen.text_to_video(
            prompt=request.prompt,
            model="sora-2",
            size=request.size,
            duration=request.duration,
            max_wait_time=600
        )
        
        if video_bytes:
            video_gen.save_video(video_bytes, output_path)
            return {
                "success": True,
                "video_id": video_id,
                "message": "Video generated successfully"
            }
        else:
            return {"success": False, "error": "Video generation failed"}
    except Exception as e:
        logging.error(f"Video generation error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/video/download/{video_id}")
async def download_video(video_id: str):
    video_path = f"/app/backend/generated_videos/{video_id}.mp4"
    if os.path.exists(video_path):
        return FileResponse(video_path, media_type="video/mp4", filename=f"{video_id}.mp4")
    raise HTTPException(status_code=404, detail="Video not found")

# Document generation endpoints
@api_router.post("/document/generate")
async def generate_document(request: DocumentGenerationRequest):
    try:
        doc_id = str(uuid.uuid4())
        os.makedirs("/app/backend/generated_docs", exist_ok=True)
        
        if request.doc_type == "pdf":
            file_path = f"/app/backend/generated_docs/{doc_id}.pdf"
            doc = SimpleDocTemplate(file_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            story.append(Paragraph(request.title, styles['Title']))
            story.append(Spacer(1, 12))
            story.append(Paragraph(request.content, styles['BodyText']))
            
            doc.build(story)
            
        elif request.doc_type == "pptx":
            file_path = f"/app/backend/generated_docs/{doc_id}.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = request.title
            content.text = request.content
            
            prs.save(file_path)
            
        elif request.doc_type == "xlsx":
            file_path = f"/app/backend/generated_docs/{doc_id}.xlsx"
            wb = Workbook()
            ws = wb.active
            ws.title = request.title[:31]  # Excel sheet name limit
            
            ws['A1'] = request.title
            ws['A2'] = request.content
            
            wb.save(file_path)
        else:
            raise HTTPException(status_code=400, detail="Invalid document type")
        
        return {
            "success": True,
            "doc_id": doc_id,
            "doc_type": request.doc_type,
            "message": "Document generated successfully"
        }
    except Exception as e:
        logging.error(f"Document generation error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/document/download/{doc_id}/{doc_type}")
async def download_document(doc_id: str, doc_type: str):
    file_path = f"/app/backend/generated_docs/{doc_id}.{doc_type}"
    if os.path.exists(file_path):
        media_types = {
            "pdf": "application/pdf",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        }
        return FileResponse(file_path, media_type=media_types.get(doc_type), filename=f"{doc_id}.{doc_type}")
    raise HTTPException(status_code=404, detail="Document not found")

# Templates endpoints
@api_router.get("/templates")
async def get_templates():
    templates = await db.templates.find({}, {"_id": 0}).to_list(1000)
    return {"templates": templates}

@api_router.get("/templates/{category}")
async def get_templates_by_category(category: str):
    templates = await db.templates.find({"category": category}, {"_id": 0}).to_list(1000)
    return {"templates": templates}

# Tutorials endpoints
@api_router.get("/tutorials")
async def get_tutorials():
    tutorials = await db.tutorials.find({}, {"_id": 0}).to_list(1000)
    return {"tutorials": tutorials}

@api_router.get("/tutorials/{category}")
async def get_tutorials_by_category(category: str):
    tutorials = await db.tutorials.find({"category": category}, {"_id": 0}).to_list(1000)
    return {"tutorials": tutorials}

# Initialize default templates and tutorials
@api_router.post("/init/data")
async def initialize_data():
    # Check if already initialized
    existing_templates = await db.templates.count_documents({})
    if existing_templates > 0:
        return {"message": "Data already initialized"}
    
    # Add default templates
    default_templates = [
        Template(
            name="Professional Blog Post Format",
            category="blog",
            description="Complete structure for writing engaging blog posts",
            content="""BLOG POST FORMAT TEMPLATE

[CATCHY TITLE - Make it attention-grabbing and SEO-friendly]

INTRODUCTION (2-3 paragraphs)
- Hook: Start with a question, statistic, or bold statement
- Context: Explain why this topic matters
- Preview: Briefly mention what you'll cover

MAIN CONTENT

Section 1: [First Main Point]
- Explanation of the concept
- Real-world example or case study
- Key takeaway in bullet form

Section 2: [Second Main Point]
- Detailed breakdown
- Supporting data or quotes
- Visual description (where image would go)

Section 3: [Third Main Point]
- Actionable advice
- Step-by-step if applicable
- Common mistakes to avoid

PRACTICAL APPLICATION
- How readers can implement this today
- Tools or resources mentioned
- Quick wins they can achieve

CONCLUSION
- Recap main points briefly
- Call-to-action (comment, share, subscribe)
- Future post teaser or related topic

AUTHOR BIO
[Your name and credentials]
[Brief description of expertise]
[Social media or website link]

---
FORMATTING TIPS:
✓ Use headings and subheadings
✓ Keep paragraphs 3-4 sentences max
✓ Include bullet points for scannability
✓ Add images every 300-400 words
✓ Aim for 1000-2000 words total"""
        ),
        Template(
            name="Business Website Structure",
            category="website",
            description="Modern responsive website layout format",
            content="""WEBSITE STRUCTURE TEMPLATE

HEADER SECTION
- Logo (top left)
- Navigation Menu: Home | About | Services | Portfolio | Contact
- CTA Button (top right): "Get Started" or "Free Consultation"

HERO SECTION (Above the fold)
- Compelling Headline (6-10 words)
- Subheadline explaining value proposition
- Primary CTA button
- Hero image or video background
- Trust badges or client logos below

FEATURES / SERVICES SECTION
Title: "What We Offer" or "Our Services"

[Service 1]
Icon/Image
Service Name
2-3 sentence description
"Learn More" link

[Service 2]
Icon/Image
Service Name
Description
Link

[Service 3]
Icon/Image
Service Name
Description
Link

ABOUT SECTION
- Company/personal story (3-4 paragraphs)
- Mission statement
- Team photo or founder image
- Stats: Years in business, Clients served, Projects completed

TESTIMONIALS / SOCIAL PROOF
Client 1: Quote + Name + Company
Client 2: Quote + Name + Company
Client 3: Quote + Name + Company

CTA SECTION
- Strong headline: "Ready to Get Started?"
- Brief text reinforcing value
- Form or button to contact

FOOTER
Column 1: Company info & logo
Column 2: Quick links
Column 3: Contact information
Column 4: Social media icons
Copyright notice at bottom

---
PAGE TYPES TO CREATE:
✓ Home (this structure)
✓ About Us
✓ Services/Products
✓ Portfolio/Case Studies
✓ Blog
✓ Contact"""
        ),
        Template(
            name="Business Presentation Format",
            category="presentation",
            description="Professional PowerPoint/Keynote structure",
            content="""PRESENTATION FORMAT TEMPLATE

SLIDE 1: TITLE SLIDE
- Presentation Title (large, bold)
- Subtitle or tagline
- Your name and title
- Company logo
- Date

SLIDE 2: AGENDA
Title: "Today's Discussion"
- Point 1: [Topic]
- Point 2: [Topic]
- Point 3: [Topic]
- Point 4: Q&A
(Keep to 4-5 items max)

SLIDE 3: PROBLEM STATEMENT
Title: "The Challenge"
- Current situation overview
- Pain points (3-4 bullets)
- Impact on business/users
- Visual: chart or image showing problem

SLIDE 4: SOLUTION OVERVIEW
Title: "Our Approach"
- High-level solution
- Key benefits (3 main points)
- Unique value proposition
- Visual: diagram or infographic

SLIDES 5-8: DETAILED BREAKDOWN
Each slide covers ONE main point:
- Clear heading
- 3-5 bullet points
- Supporting visual
- Keep text minimal

SLIDE 9: CASE STUDY / PROOF
Title: "Success Story" or "Results"
- Before/After comparison
- Key metrics and improvements
- Client quote (if applicable)
- Visual evidence

SLIDE 10: NEXT STEPS
Title: "Moving Forward"
- Action items
- Timeline
- Resources needed
- Decision points

SLIDE 11: CALL TO ACTION
- Clear next step
- Contact information
- Meeting follow-up details

SLIDE 12: THANK YOU
- "Thank You" or "Questions?"
- Your contact details
- Company website/social media
- Offer to send presentation deck

---
DESIGN PRINCIPLES:
✓ One idea per slide
✓ Max 6 bullets per slide
✓ Use high-quality images
✓ Consistent fonts (2 max)
✓ Brand colors throughout
✓ White space is good!
✓ 30-60 second per slide rule"""
        ),
        Template(
            name="Data Spreadsheet Format",
            category="website",
            description="Excel/Sheets organization template",
            content="""SPREADSHEET FORMAT TEMPLATE

SHEET 1: SUMMARY DASHBOARD
Row 1: Title and Date
Row 3-5: Key Metrics Summary
- Total Revenue
- Total Expenses  
- Net Profit
- Growth %

Row 7: Monthly Comparison Chart Reference
Row 9-20: Quick Stats Table

SHEET 2: DETAILED DATA
Column Headers (Row 1):
A: Date
B: Category
C: Description
D: Amount
E: Payment Method
F: Status
G: Notes

Data Rows:
- One transaction per row
- Consistent date format
- Clear categorization
- Numerical data right-aligned

Row (Last): TOTALS
Sum formulas for numerical columns

SHEET 3: MONTHLY BREAKDOWN
Table Structure:
         Jan  Feb  Mar  Apr  May  Jun  Jul  Aug  Sep  Oct  Nov  Dec  TOTAL
Income   XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  =SUM
Expense  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  =SUM
Profit   XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  XXX  =SUM

SHEET 4: CHARTS & VISUALS
- Revenue Trend (Line Chart)
- Expense Breakdown (Pie Chart)
- Category Comparison (Bar Chart)
- Year-over-Year Growth

---
FORMATTING BEST PRACTICES:
✓ Freeze top row and first column
✓ Use consistent number formats
✓ Color-code headers
✓ Apply data validation
✓ Add filters to header row
✓ Protect formula cells
✓ Document formulas in comments
✓ Create named ranges for key data"""
        )
    ]
    
    for template in default_templates:
        await db.templates.insert_one(template.model_dump())
    
    # Add default tutorials
    default_tutorials = [
        Tutorial(
            title="Getting Started with P07 Chat",
            description="Master the art of AI conversations with P07's dual-model chat system",
            content="""HOW TO USE P07 CHAT EFFECTIVELY:

1. SELECT YOUR AI MODEL: Choose between GPT-5.2 (for creative and detailed responses) or Claude Sonnet 4.5 (for analytical and precise answers) using the model selector buttons at the top.

2. ASK CLEAR QUESTIONS: For best results, be specific with your queries. Instead of "tell me about physics," ask "explain Newton's laws of motion with real-world examples."

3. REQUEST FORMATTED NOTES: When you need study materials or documentation, ask P07 to "give me notes on [topic]" and it will provide structured, formatted educational content with headings, bullet points, and clear explanations.

4. UTILIZE CONVERSATION HISTORY: P07 remembers your chat history within each session, so you can build upon previous questions and have natural, flowing conversations.

5. SWITCH MODELS FOR DIFFERENT TASKS: Use GPT-5.2 for creative writing, brainstorming, and detailed explanations. Use Claude 4.5 for code review, logical analysis, and fact-checking.

TIPS FOR MAXIMUM PRODUCTIVITY:
- Start new sessions for different topics to keep conversations organized
- Ask follow-up questions to dig deeper into complex subjects
- Request examples, analogies, or case studies for better understanding
- Use P07 for summarizing long documents or articles
- Ask for step-by-step guides when learning new skills""",
            category="chat"
        ),
        Tutorial(
            title="Creating Images with SP07",
            description="Generate stunning, high-resolution images using Gemini Nano Banana",
            content="""HOW TO USE SP07 IMAGE GENERATION:

1. WRITE DESCRIPTIVE PROMPTS: The more detailed your description, the better your image. Instead of "a cat," write "a majestic Persian cat with golden fur, sitting on a velvet cushion in a Victorian-style room with soft natural lighting."

2. SPECIFY STYLE AND MOOD: Include artistic style in your prompt - "photorealistic," "oil painting style," "minimalist digital art," "cinematic," "vintage photograph," etc.

3. INCLUDE COMPOSITION DETAILS: Mention camera angles ("wide-angle shot," "close-up portrait"), lighting ("golden hour," "dramatic shadows"), and perspective to guide the generation.

4. BE PATIENT WITH GENERATION: Image creation takes a few moments. Wait for the complete generation - the result will be worth it with ultra-high resolution and realistic details.

5. DOWNLOAD AND REFINE: Once generated, download your image. If not perfect, create a new prompt with adjusted details based on what you learned.

EXAMPLE PROMPTS THAT WORK WELL:
- "A futuristic city skyline at sunset, with flying cars and neon lights, cyberpunk style, highly detailed"
- "Professional headshot of a business executive, studio lighting, neutral background, photorealistic"
- "Abstract geometric pattern in gold and navy blue, luxury design, 4K quality"
- "Cozy coffee shop interior, warm lighting, wooden furniture, people working on laptops, wide angle"

AVOID:
- Vague prompts like "make something cool"
- Requesting copyrighted characters or logos
- Too many conflicting style directions in one prompt""",
            category="image"
        ),
        Tutorial(
            title="Document Generation with FP07",
            description="Create professional PDFs, PowerPoints, and Excel files instantly",
            content="""HOW TO USE FP07 DOCUMENT GENERATOR:

1. CHOOSE YOUR DOCUMENT TYPE FIRST: Before writing content, select whether you need a PDF (for reports/documents), PPTX (for presentations), or XLSX (for spreadsheets/data tables) to format your content appropriately.

2. CRAFT A CLEAR TITLE: Your document title should be concise and descriptive. This becomes the filename and the main heading. Examples: "Q4 2026 Marketing Report," "Product Launch Presentation," "Monthly Sales Data."

3. STRUCTURE YOUR CONTENT PROPERLY: 
   - For PDFs: Write in paragraphs with clear sections
   - For PPTX: Use bullet points and short phrases for each slide point
   - For XLSX: Organize data in rows and columns format

4. GENERATE AND REVIEW: Click Generate, wait a few seconds, then preview your document. FP07 formats it professionally with proper styling and layout.

5. DOWNLOAD IN THE RIGHT FORMAT: Use the download button to save your document. Open it in the appropriate software (PDF Reader, PowerPoint, Excel) to verify and make final adjustments if needed.

CONTENT FORMATTING TIPS:
- Use line breaks to separate sections
- Start main points with capital letters
- For presentations: Keep each point under 15 words
- For spreadsheets: Use commas to separate column data

BEST PRACTICES:
- Keep titles under 50 characters
- Break long content into sections with headings
- Use professional language for business documents
- Include dates and version numbers in titles
- Proofread before generating to avoid regeneration""",
            category="documents"
        ),
        Tutorial(
            title="Mastering ChatGPT & Gemini for Work",
            description="Learn professional AI tools usage strategies from P07 experts",
            content="""HOW TO USE CHATGPT & GEMINI EFFICIENTLY AT WORK:

1. UNDERSTAND EACH AI'S STRENGTHS: ChatGPT excels at creative writing, coding, and explanations. Gemini is better for research, fact-checking, and Google-integrated tasks. Use the right tool for the right job.

2. CREATE EFFECTIVE PROMPTS: Start with context ("I am a marketing manager working on..."), then state your request clearly ("create a 30-day social media campaign for..."), and specify format ("in table format with dates and post ideas").

3. USE CONVERSATION CHAINS: Break complex tasks into steps. Start broad, then narrow down with follow-up questions. Example: First ask for strategy overview, then request specific tactics, finally get implementation details.

4. MAINTAIN CONVERSATION CONTEXT: Both AIs remember your chat history. Reference previous answers with "based on what you just said" or "expand on point 3" to build comprehensive solutions.

5. LEVERAGE FOR DAILY WORKFLOWS: Use AI for email drafting, meeting notes summarization, report writing, data analysis, brainstorming, research, and learning new skills. Integrate into your routine tasks.

WORKPLACE APPLICATIONS:
- Email responses: "Write a professional email declining this meeting politely"
- Reports: "Summarize these sales figures into a 2-paragraph executive summary"
- Learning: "Explain blockchain technology as if I'm explaining to my CEO"
- Analysis: "What are the pros and cons of this marketing strategy?"
- Content: "Create 10 LinkedIn post ideas about [your industry]"

PRODUCTIVITY HACKS:
- Save your best prompts as templates
- Use AI to critique your own work before submitting
- Ask for multiple versions to choose from
- Get AI to simplify complex jargon for presentations""",
            category="chat"
        )
    ]
    
    for tutorial in default_tutorials:
        await db.tutorials.insert_one(tutorial.model_dump())
    
    return {"message": "Data initialized successfully"}

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()